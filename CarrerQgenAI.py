import streamlit as st
from docx import Document
import PyPDF2
import pptx
import json
from io import BytesIO
import re
import requests

# Define the API endpoint and access token
API_URL = "https://api.psnext.info/api/chat"
PSCHATACCESSTOKEN = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJVc2VySW5mbyI6eyJpZCI6MzcxMzcsInJvbGVzIjpbImRlZmF1bHQiXSwicGF0aWQiOiJmYTRlMzZhNC1mMWU5LTRjMjktODYwMi0wZDU1NGFmMGIxYzcifSwiaWF0IjoxNzI2MDUwNDI2LCJleHAiOjE3Mjg2NDI0MjZ9.gw74MhUO6rO3wrauMHQxdm8PWK6RBJAh6v7yIFSS8zA"

# Function to extract text from a Word document
def extract_text_from_word(docx_file):
    doc = Document(docx_file)
    return '\n'.join([para.text for para in doc.paragraphs])

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    pdf_text = ""
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        pdf_text += page.extract_text()
    return pdf_text

# Function to extract text from a PPT file
def extract_text_from_ppt(ppt_file):
    ppt = pptx.Presentation(ppt_file)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Function to extract text from different file types
def extract_text_from_file(file):
    if file.type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_word(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_ppt(file)
    else:
        return ""

# Function to call PSNext API for CV matching and rating
def get_cv_match(cv_text, job_description):
    payload = {
        "message": f"Evaluate this CV against the following job description, providing a rating out of 10 and feedback:\n\nJob Description:\n{job_description}\n\nCV:\n{cv_text}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                result = message.get('content', 'No content returned from the API.')
                match = re.search(r'Rating: (\d+)/10', result)
                if match:
                    rating = int(match.group(1))
                    feedback = result.split('\n', 1)[1] if '\n' in result else ''
                    return rating, feedback
                else:
                    return 0, 'Unable to extract rating from the response.'
        return 0, 'No assistant message found in the API response.'
    else:
        return 0, f"Error: {response.status_code}, {response.text}"

# Function to generate case study questions
def generate_case_study_questions(job_description, years_of_experience, industry, difficulty_level):
    payload = {
        "message": f"Generate a set of case study questions based on the following job description, {years_of_experience} years of experience in the {industry} industry, with a difficulty level of {difficulty_level}:\n\nJob Description:\n{job_description}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                return message.get('content', 'No content returned from the API.')
        return 'No assistant message found in the API response.'
    else:
        return f"Error: {response.status_code}, {response.text}"

# Function to match case study answers
def match_case_study_answers(question, provided_answer):
    payload = {
        "message": f"Evaluate the following case study question and answer, providing a rating out of 10 and a brief feedback:\n\nQuestion:\n{question}\n\nAnswer:\n{provided_answer}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                result = message.get('content', 'No content returned from the API.')
                match = re.search(r'Rating: (\d+)/10', result)
                if match:
                    rating = int(match.group(1))
                    feedback = result.split('\n', 1)[1] if '\n' in result else ''
                    return rating, feedback
                else:
                    return 0, 'Unable to extract rating from the response.'
        return 0, 'No assistant message found in the API response.'
    else:
        return 0, f"Error: {response.status_code}, {response.text}"

# Function to compare two sets of texts (questions and solutions)
def compare_question_solution(question_text, solution_text):
    payload = {
        "message": f"Compare the following question document with the solution document. Provide feedback, rating out of 10, and suggestions for improvement:\n\nQuestion Document:\n{question_text}\n\nSolution Document:\n{solution_text}",
        "options": {"model": "gpt35turbo"}
    }
    
    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }
    
    response = requests.post(API_URL, headers=headers, data=json.dumps(payload))
    
    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                result = message.get('content', 'No content returned from the API.')
                match = re.search(r'Rating: (\d+)/10', result)
                if match:
                    rating = int(match.group(1))
                    feedback = result.split('\n', 1)[1] if '\n' in result else ''
                    return rating, feedback
                else:
                    return 0, 'Unable to extract rating from the response.'
        return 0, 'No assistant message found in the API response.'
    else:
        return 0, f"Error: {response.status_code}, {response.text}"

# Emoji mapping based on rating
def get_rating_emoji(rating):
    if rating >= 9:
        return "🌟 Excellent"
    elif rating >= 7:
        return "👍 Good"
    elif rating >= 5:
        return "👌 Average"
    elif rating >= 3:
        return "🤔 Okay"
    else:
        return "👎 Bad"

# Main app page
def main_app():
    st.title("CareerQgen-Your AI-Powered Staffing Solution")
    st.subheader("Optimize Your Hiring Process with AI")

    tabs = st.tabs(["CV Matching", "Case Study Generation", "Case Study Evaluation", "Document Comparison"])

    # Tab 1: CV Matching
    with tabs[0]:
        st.header("CV Matching")
        st.write("Upload a CV (PDF or Word) and provide a job description to get a match evaluation.")
        
        uploaded_cv = st.file_uploader("Upload Candidate CV (PDF or Word)", type=["pdf", "docx"], key="cv_upload")
        job_description = st.text_area("Job Description", height=200, key="job_description_cv")
        
        result_output = st.empty()

        if uploaded_cv and job_description:
            if uploaded_cv.type == "application/pdf":
                cv_text = extract_text_from_pdf(uploaded_cv)
            elif uploaded_cv.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                cv_text = extract_text_from_word(uploaded_cv)

            if cv_text.strip():
                if st.button("Get CV Match", key="cv_match_button"):
                    with st.spinner("Processing..."):
                        rating, feedback = get_cv_match(cv_text, job_description)
                        result_output.text_area("Match Feedback", feedback, height=200, key="cv_feedback")
                        st.write(f"Rating: {rating}/10 {get_rating_emoji(rating)}")
            else:
                st.error("No text could be extracted from the uploaded CV.")

    # Tab 2: Case Study Generation
    with tabs[1]:
        st.header("Case Study Generation")
        st.write("Generate case study questions based on a job description.")
        
        job_description = st.text_area("Job Description", height=200, key="job_description_case_study")
        years_of_experience = st.number_input("Years of Experience", min_value=0, max_value=50, step=1, key="experience")
        industry = st.selectbox("Industry/Domain", ["Finance", "Supply Chain", "Technology", "Healthcare", "Education", "Other"], key="industry")
        difficulty_level = st.selectbox("Difficulty Level", ["Easy", "Moderate", "Difficult"], key="difficulty")

        case_study_output = st.empty()

        if job_description and years_of_experience:
            if st.button("Generate Case Study Questions", key="generate_case_study_button"):
                with st.spinner("Processing..."):
                    case_study = generate_case_study_questions(job_description, years_of_experience, industry, difficulty_level)
                    case_study_output.text_area("Generated Case Study Questions", case_study, height=200, key="case_study_output")

    # Tab 3: Case Study Answer Matching
    with tabs[2]:
        st.header("Case Study Answer Matching")
        st.write("Provide a case study question and the candidate's answer to get a match evaluation.")
        case_study_question = st.text_area("Case Study Question", height=200, key="case_study_question")
        candidate_answer = st.text_area("Candidate's Answer", height=200, key="candidate_answer")

        result_output = st.empty()

        if case_study_question and candidate_answer:
            if st.button("Evaluate Answer", key="evaluate_answer_button"):
                with st.spinner("Processing..."):
                    rating, feedback = match_case_study_answers(case_study_question, candidate_answer)
                    result_output.text_area("Answer Feedback", feedback, height=200, key="answer_feedback")
                    st.write(f"Rating: {rating}/10 {get_rating_emoji(rating)}")

    # Tab 4: Document Comparison
    with tabs[3]:
        st.header("Document Comparison")
        st.write("Upload two documents (e.g., a question and solution document) to compare and get feedback.")

        uploaded_question_doc = st.file_uploader("Upload Question Document (PDF, Word, or PPT)", type=["pdf", "docx", "pptx"], key="question_doc")
        uploaded_solution_doc = st.file_uploader("Upload Solution Document (PDF, Word, or PPT)", type=["pdf", "docx", "pptx"], key="solution_doc")

        result_output = st.empty()

        if uploaded_question_doc and uploaded_solution_doc:
            question_text = extract_text_from_file(uploaded_question_doc)
            solution_text = extract_text_from_file(uploaded_solution_doc)

            if st.button("Compare Documents", key="compare_documents_button"):
                with st.spinner("Processing..."):
                    rating, feedback = compare_question_solution(question_text, solution_text)
                    result_output.text_area("Comparison Feedback", feedback, height=200, key="comparison_feedback")
                    st.write(f"Rating: {rating}/10 {get_rating_emoji(rating)}")

# Run the app
if __name__ == "__main__":
    main_app()
